import streamlit as st
import pandas as pd
import sqlite3
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO

# Importações para Exportações de Relatórios
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt

# ==========================================
# CONFIGURAÇÃO DA INTERFACE
# ==========================================
st.set_page_config(page_title="Sistema Avançado de Inteligência de Ativos - DRC", layout="wide")

DB_NAME = "ativos_engenharia_v2.db"

# ==========================================
# REQUISITO 1, 6 & 7: ESTRUTURA DE BANCO DE DADOS RELACIONAL
# ==========================================
def init_db():
    conn = sqlite3.connect(DB_NAME)
    cursor = conn.cursor()
    
    # 1) Tabela de UPs (Cadastro, Edição e Exclusão)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS tabela_up (
        codigo_up TEXT PRIMARY KEY,
        nomenclatura_up TEXT NOT NULL,
        tipo_massa TEXT DEFAULT 'Individual'
    )""")
    
    # Bases de Preços Referenciais carregadas (com suporte a i0 diferentes)
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS base_referencial (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        banco TEXT,         -- SABESP, SINAPI, TCPO
        codigo_item TEXT,
        descricao TEXT,
        unidade TEXT,
        disciplina TEXT,    -- Civil, Elétrica, Hidromecânica
        i0_data TEXT,       -- Formato YYYY-MM
        preco REAL
    )""")

    # 5) Tabela de Índices IPCA Históricos por i0
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS indices_ipca (
        i0_data TEXT PRIMARY KEY,
        fator_acumulado REAL
    )""")
    
    # 6) Tabela de Propostas de Fornecedores e Contratos
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS propostas_fornecedores (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        tipo_serviço_ativo TEXT,
        codigo_up TEXT,
        preco_fornecedor REAL,
        item_similar_banco TEXT,
        preco_contratado REAL,
        FOREIGN KEY (codigo_up) REFERENCES tabela_up(codigo_up)
    )""")
    
    # Carga inicial padrão para Forward Momentum
    cursor.execute("SELECT COUNT(*) FROM tabela_up")
    if cursor.fetchone()[0] == 0:
        cursor.executemany("INSERT INTO tabela_up VALUES (?,?,?)", [
            ('01', 'Instalações Civis de Saneamento', 'Individual'),
            ('02', 'Sistemas Elétricos e Subestações', 'Individual'),
            ('03', 'Conjuntos Hidromecânicos e Bombas', 'Individual')
        ])
        
        # Inserindo dados de i0 diferentes para cálculo de inflação (Janeiro a Junho de 2026)
        cursor.executemany("INSERT INTO base_referencial (banco, codigo_item, descricao, unidade, disciplina, i0_data, preco) VALUES (?,?,?,?,?,?,?)", [
            ('SABESP', 'S01', 'Tubo PVC Hidromecânico', 'm', 'Hidromecânica', '2026-01', 100.0),
            ('SABESP', 'S01', 'Tubo PVC Hidromecânico', 'm', 'Hidromecânica', '2026-06', 112.0),
            ('SINAPI', 'SI02', 'Cabo Cobre Isolado', 'm', 'Elétrica', '2026-01', 50.0),
            ('SINAPI', 'SI02', 'Cabo Cobre Isolado', 'm', 'Elétrica', '2026-06', 54.5),
            ('TCPO', 'T03', 'Concreto Armado Estrutural', 'm3', 'Civil', '2026-01', 400.0),
            ('TCPO', 'T03', 'Concreto Armado Estrutural', 'm3', 'Civil', '2026-06', 424.0),
        ])
        
        cursor.executemany("INSERT INTO indices_ipca VALUES (?,?)", [
            ('2026-01', 1.0000), ('2026-03', 1.0120), ('2026-04', 1.0185), ('2026-06', 1.0310)
        ])
        
        cursor.executemany("INSERT INTO propostas_fornecedores (tipo_serviço_ativo, codigo_up, preco_fornecedor, item_similar_banco, preco_contratado) VALUES (?,?,?,?,?)", [
            ('Bomba de recalque auxiliar', '03', 113.0, 'S01', 111.0),
            ('Reforma de painel elétrico', '02', 62.0, 'SI02', 53.0),
            ('Base de concreto para motor', '01', 390.0, 'T03', 405.0)
        ])
        
    conn.commit()
    conn.close()

init_db()

def get_db():
    return sqlite3.connect(DB_NAME)

# ==========================================
# INTERFACE PRINCIPAL - TÍTULO E MENU
# ==========================================
st.title("⚡ Sistema de Inteligência de Ativos, DRC Regulatório e Cadeia de Preços")
menu = st.sidebar.selectbox("Módulos do Sistema", [
    "Gerenciar UPs (CRUD)", 
    "Contratos & Fornecedores", 
    "Cálculo de Inflação (Setorial e IPCA)", 
    "Dashboard & Relatórios"
])

# ==========================================
# MODULO 1: GERENCIAR UPs (CADASTRO, EDIÇÃO E EXCLUSÃO)
# ==========================================
if menu == "Gerenciar UPs (CRUD)":
    st.header("⚙️ Cadastro, Edição e Exclusão de Unidades de Patrimônio (UP)")
    
    conn = get_db()
    
    # Formulário de Cadastro/Edição
    tipo_acao = st.radio("Escolha a operação", ["Cadastrar Nova UP", "Editar UP Existente", "Excluir UP"])
    
    if tipo_acao == "Cadastrar Nova UP":
        with st.form("cadastrar_up"):
            cod = st.text_input("Código da UP (Ex: 04)")
            nome = st.text_input("Nomenclatura da UP")
            massa = st.selectbox("Tipo de Massa", ["Individual", "Global"])
            sub = st.form_submit_with_button("Salvar UP")
            if sub and cod and nome:
                try:
                    conn.execute("INSERT INTO tabela_up VALUES (?,?,?)", (cod, nome, massa))
                    conn.commit()
                    st.success(f"UP {cod} cadastrada com sucesso!")
                except sqlite3.IntegrityError:
                    st.error("Código de UP já existente.")
                    
    elif tipo_acao == "Editar UP Existente":
        ups = pd.read_sql_query("SELECT * FROM tabela_up", conn)
        if not ups.empty:
            cod_selecionado = st.selectbox("Selecione a UP para editar", ups['codigo_up'])
            up_atual = ups[ups['codigo_up'] == cod_selecionado].iloc[0]
            
            with st.form("editar_up"):
                novo_nome = st.text_input("Nova Nomenclatura", value=up_atual['nomenclatura_up'])
                novo_massa = st.selectbox("Novo Tipo de Massa", ["Individual", "Global"], index=0 if up_atual['tipo_massa'] == 'Individual' else 1)
                sub_edit = st.form_submit_with_button("Atualizar UP")
                if sub_edit:
                    conn.execute("UPDATE tabela_up SET nomenclatura_up = ?, tipo_massa = ? WHERE codigo_up = ?", (novo_nome, novo_massa, cod_selecionado))
                    conn.commit()
                    st.success("UP atualizada com sucesso!")
        else:
            st.warning("Nenhuma UP cadastrada.")
            
    elif tipo_acao == "Excluir UP":
        ups = pd.read_sql_query("SELECT * FROM tabela_up", conn)
        if not ups.empty:
            cod_deletar = st.selectbox("Selecione a UP para remover", ups['codigo_up'])
            if st.button("Confirmar Exclusão Definitiva"):
                conn.execute("DELETE FROM tabela_up WHERE codigo_up = ?", (cod_deletar,))
                conn.commit()
                st.success("UP removida do banco de dados com sucesso!")
                st.rerun()
                
    st.subheader("Lista Atual de UPs no Sistema")
    df_lista = pd.read_sql_query("SELECT * FROM tabela_up", conn)
    st.dataframe(df_lista, use_container_width=True)
    conn.close()

# ==========================================
# MODULO 2: CONTRATOS & FORNECEdoRES (REQUISITOS 6, 7 & 8)
# ==========================================
elif menu == "Contratos & Fornecedores":
    st.header("🤝 Gestão de Propostas de Fornecedores e Análise DRC Parametrizável")
    
    # REQUISITO 7: Parâmetros do DRC na tela
    st.sidebar.markdown("### Configurações do DRC")
    limite_drc = st.sidebar.slider("Margem de Tolerância DRC (%)", 1.0, 15.0, 5.0, step=0.5)
    
    conn = get_db()
    
    # 6) CRUD Preço Fornecedor e Contratado com Classificação Automática por UP
    st.subheader("Adicionar/Editar Registro de Proposta Comercial")
    with st.form("form_fornecedor"):
        col1, col2 = st.columns(2)
        with col1:
            ativo_servico = st.text_input("Tipo de Ativo ou Serviço do Fornecedor")
            p_forn = st.number_input("Preço Proposto do Fornecedor (R$)", min_value=0.0, value=100.0)
            p_cont = st.number_input("Preço Final Contratado (R$)", min_value=0.0, value=95.0)
        with col2:
            # Seleção de itens similares dos bancos para cruzamento automático
            df_itens = pd.read_sql_query("SELECT codigo_item, banco, descricao, preco FROM base_referencial", conn)
            item_sim_opcoes = [f"{r['codigo_item']} - {r['banco']} - {r['descricao']} (Ref: R$ {r['preco']})" for _, r in df_itens.iterrows()]
            item_selecionado = st.selectbox("Selecione o Item de Referência Similar (Sabesp/SINAPI/TCPO)", item_sim_opcoes) if item_sim_opcoes else "Nenhum cadastrado"
            
            # Vinculação automática com a UP baseada na tabela
            df_ups = pd.read_sql_query("SELECT codigo_up, nomenclatura_up FROM tabela_up", conn)
            up_opcoes = [f"{r['codigo_up']} - {r['nomenclatura_up']}" for _, r in df_ups.iterrows()]
            up_selecionada = st.selectbox("Classificação Automática de UP Destino", up_opcoes)
            
        submit_forn = st.form_submit_with_button("Salvar Registro Comercial")
        if submit_forn and item_selecionado != "Nenhum cadastrado":
            cod_item_clean = item_selecionado.split(" - ")[0]
            cod_up_clean = up_selecionada.split(" - ")[0]
            conn.execute("""
                INSERT INTO propostas_fornecedores (tipo_serviço_ativo, codigo_up, preco_fornecedor, item_similar_banco, preco_contratado) 
                VALUES (?, ?, ?, ?, ?)
            """, (ativo_servico, cod_up_clean, p_forn, cod_item_clean, p_cont))
            conn.commit()
            st.success("Transação comercial inserida e classificada por UP automaticamente!")

    # Exibição da Análise Comparativa com as Cores de Alerta (REQUISITO 8)
    st.subheader("Análise de Conformidade DRC e Sinais de Alerta")
    
    query_analise = """
        SELECT pf.id, pf.tipo_serviço_ativo, pf.codigo_up, tu.nomenclatura_up, pf.preco_fornecedor, pf.preco_contratado,
               br.banco, br.preco as preco_referencia_banco, br.disciplina
        FROM propostas_fornecedores pf
        JOIN tabela_up tu ON pf.codigo_up = tu.codigo_up
        JOIN base_referencial br ON pf.item_similar_banco = br.codigo_item
        GROUP BY pf.id
    """
    df_analise = pd.read_sql_query(query_analise, conn)
    
    if not df_analise.empty:
        # Funções para cálculo de desvio e Alerta Visual
        def avaliar_drc(row):
            ref = row['preco_referencia_banco']
            forn = row['preco_fornecedor']
            desvio = ((forn - ref) / ref) * 100
            
            # REQUISITO 8: Lógica de Alertas por cores
            if abs(desvio) <= limite_drc:
                return "🟢 Conforme", f"background-color: #d4edda; color: #155724;"
            elif abs(desvio) <= (limite_drc + 5.0):
                return "🟡 Atenção", f"background-color: #fff3cd; color: #856404;"
            else:
                return "🔴 Não Conforme", f"background-color: #f8d7da; color: #721c24;"

        alertas_status = []
        estilos_linhas = []
        desvios_valores = []
        
        for idx, row in df_analise.iterrows():
            status, estilo = avaliar_drc(row)
            alertas_status.append(status)
            desvios_valores.append(((row['preco_fornecedor'] - row['preco_referencia_banco']) / row['preco_referencia_banco']) * 100)
            
        df_analise['Desvio (%)'] = desvios_valores
        df_analise['Status DRC'] = alertas_status
        
        # Formatando tabela estilizada para exibição em tela
        st.dataframe(df_analise.style.apply(lambda x: [avaliar_drc(row)[1] if x.name == 'Status DRC' else '' for row in df_analise.itertuples()], axis=0), use_container_width=True)
        
        # Botão para Limpeza/Exclusão de registros comerciais para controle de dados
        st.markdown("---")
        id_deletar = st.selectbox("Selecione o ID do registro comercial para excluir se necessário:", df_analise['id'])
        if st.button("Excluir Registro de Preço Selecionado"):
            conn.execute("DELETE FROM propostas_fornecedores WHERE id = ?", (id_deletar,))
            conn.commit()
            st.success("Registro excluído!")
            st.rerun()
    else:
        st.info("Nenhum dado comercial de fornecedores disponível.")
    conn.close()

# ==========================================
# MODULO 3: CÁLCULO DE INFLAÇÃO SETORIAL E IPCA (REQUISITOS 4 & 5)
# ==========================================
elif menu == "Cálculo de Inflação (Setorial e IPCA)":
    st.header("🏭 Cálculo de Inflação de Preços por Disciplina e Indexadores Históricos")
    
    conn = get_db()
    
    # 4) Inflação de Preços por Disciplina Comparando Diferentes i0 do próprio banco
    st.subheader("Cálculo 1: Inflação Setorial Interna por Banco ($i_0$ Base vs $i_0$ Destino)")
    
    col_inf1, col_inf2 = st.columns(2)
    with col_inf1:
        banco_sel = st.selectbox("Escolha o Banco de Dados", ["SABESP", "SINAPI", "TCPO"])
    with col_inf2:
        disciplina_sel = st.selectbox("Selecione a Disciplina Técnica", ["Civil", "Elétrica", "Hidromecânica"])
        
    df_i0_disponiveis = pd.read_sql_query(f"SELECT DISTINCT i0_data FROM base_referencial WHERE banco='{banco_sel}' AND disciplina='{disciplina_sel}'", conn)
    
    if len(df_i0_disponiveis) >= 2:
        list_i0 = df_i0_disponiveis['i0_data'].tolist()
        i0_inicial = st.selectbox("i0 Inicial (Base)", list_i0, index=0)
        i0_final = st.selectbox("i0 Final (Destino)", list_i0, index=len(list_i0)-1)
        
        if st.button("Calcular Inflação Setorial"):
            p_ini = pd.read_sql_query(f"SELECT AVG(preco) as m FROM base_referencial WHERE banco='{banco_sel}' AND disciplina='{disciplina_sel}' AND i0_data='{i0_inicial}'", conn)['m'].values[0]
            p_fim = pd.read_sql_query(f"SELECT AVG(preco) as m FROM base_referencial WHERE banco='{banco_sel}' AND disciplina='{disciplina_sel}' AND i0_data='{i0_final}'", conn)['m'].values[0]
            
            if p_ini and p_fim:
                variacao_setorial = ((p_fim - p_ini) / p_ini) * 100
                st.metric(label=f"Inflação Acumulada no Banco {banco_sel} ({disciplina_sel})", value=f"{variacao_setorial:.2f}%", delta=f"{p_fim - p_ini:.2f} R$ em relação à base")
    else:
        st.warning("É necessário fazer upload ou ter cadastrado pelo menos duas datas de i0 diferentes para este banco e disciplina para calcular variações.")

    # 5) Inflação por IPCA em relação aos i0 carregados pelo usuário
    st.write("---")
    st.subheader("Cálculo 2: Reajuste Real vs Atualização Monetária (IPCA)")
    df_ipca = pd.read_sql_query("SELECT * FROM indices_ipca", conn)
    st.write("Fatores do IPCA disponíveis carregados pelo usuário (Ano Referência 2026):")
    st.dataframe(df_ipca.T)
    
    df_itens_geral = pd.read_sql_query("SELECT id, banco, codigo_item, descricao, preco, i0_data FROM base_referencial", conn)
    if not df_itens_geral.empty and not df_ipca.empty:
        item_ipca_sel = st.selectbox("Selecione o material/equipamento para analisar o IPCA", [f"{r['id']} - {r['banco']} - {r['descricao']} ({r['i0_data']})" for _, r in df_itens_geral.iterrows()])
        id_item_clean = item_ipca_sel.split(" - ")[0]
        
        linha_item = df_itens_geral[df_itens_geral['id'] == int(id_item_clean)].iloc[0]
        fator_base = df_ipca[df_ipca['i0_data'] == linha_item['i0_data']]['fator_acumulado'].values
        
        target_i0 = st.selectbox("Selecione o i0 do IPCA de destino para simulação de reajuste", df_ipca['i0_data'].unique())
        fator_destino = df_ipca[df_ipca['i0_data'] == target_i0]['fator_acumulado'].values[0]
        
        if len(fator_base) > 0:
            preco_corrigido_ipca = linha_item['preco'] * (fator_destino / fator_base[0])
            st.success(f"O preço original de R$ {linha_item['preco']:.2f} em {linha_item['i0_data']}, corrigido pelo IPCA para {target_i0}, equivale a **R$ {preco_corrigido_ipca:.2f}**.")
        else:
            st.error("Não há fator IPCA cadastrado para o mês base deste item.")
            
    conn.close()

# ==========================================
# MODULO 4: DASHBOARDS, GRÁFICOS & EXPORTAÇÕES (REQUISITOS 9 & 10)
# ==========================================
elif menu == "Dashboard & Relatórios":
    st.header("📊 Inteligência Visual e Exportação Multiformato de Relatórios Técnicos")
    
    conn = get_db()
    df_propostas = pd.read_sql_query("SELECT * FROM propostas_fornecedores", conn)
    df_referencias = pd.read_sql_query("SELECT * FROM base_referencial", conn)
    conn.close()
    
    if not df_propostas.empty and not df_referencias.empty:
        # Mapeamento e união de bases de dados para consolidar gráficos
        df_total = df_propostas.merge(df_referencias, left_on='item_similar_banco', right_on='codigo_item')
        
        # REQUISITO 9: Geração de 4 Tipos de Gráficos Distintos Exigidos
        st.subheader("Visualizações Estratégicas Integradas")
        g1, g2 = st.columns(2)
        g3, g4 = st.columns(2)
        
        with g1:
            st.markdown("**1. Gráfico de Linha (Série Temporal por Banco)**")
            fig_lin = px.line(df_referencias, x='i0_data', y='preco', color='banco', markers=True, title="Evolução Histórica das Bases de Engenharia")
            st.plotly_chart(fig_lin, use_container_width=True)
            
        with g2:
            st.markdown("**2. Gráfico de Barras (Comparativo de Preços por Ativo)**")
            fig_bar = px.bar(df_total, x='tipo_serviço_ativo', y=['preco_fornecedor', 'preco_contratado', 'preco'], barmode='group', title="Fornecedor vs Contratado vs Banco de Dados")
            st.plotly_chart(fig_bar, use_container_width=True)
            
        with g3:
            st.markdown("**3. Heatmap de Correlação / Intensidade de Preços**")
            # Montando matriz pivotada simples para simular intensidade de preços por banco/disciplina
            df_pivot = df_referencias.pivot_table(index='disciplina', columns='banco', values='preco', aggfunc='mean').fillna(0)
            fig_heat = px.imshow(df_pivot, text_auto=True, color_continuous_scale='YlOrRd', title="Concentração de Custos Médios por Categoria")
            st.plotly_chart(fig_heat, use_container_width=True)
            
        with g4:
            st.markdown("**4. Gráfico Bar Incremental (Cadeia de Valor Acumulado)**")
            fig_inc = go.Figure(go.Waterfall(
                name = "Ajustes", orientation = "v",
                measure = ["relative", "relative", "total"],
                x = ["Preço Médio Referência", "Margem Fornecedores", "Custo Total Contratado Estimado"],
                y = [df_total['preco'].mean(), (df_total['preco_fornecedor'].mean() - df_total['preco'].mean()), df_total['preco_contratado'].mean()],
                connector = {"line":{"color":"rgb(63, 63, 63)"}},
            ))
            fig_inc.update_layout(title = "Cascata Incremental de Desvios de Mercado")
            st.plotly_chart(fig_inc, use_container_width=True)

        # REQUISITO 10: MOTOR DE EXPORTAÇÃO DE RELATÓRIOS MULTIFORMATO
        st.write("---")
        st.subheader("💾 Exportação Oficial de Relatórios de Auditoria")
        st.write("Selecione os formatos abaixo para fazer o download dos relatórios consolidados gerados em tempo real:")
        
        col_exp1, col_exp2, col_exp3 = st.columns(3)
        
        # 10.1) GERAÇÃO DO EXCEL AVANÇADO
        with col_exp1:
            output_excel = BytesIO()
            with pd.ExcelWriter(output_excel, engine='xlsxwriter') as writer:
                df_total.to_excel(writer, sheet_name='Analise_DRC', index=False)
            st.download_button(
                label="📥 Baixar em Excel (.xlsx)",
                data=output_excel.getvalue(),
                file_name="Relatorio_Auditoria_Ativos.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            )
            
        # 10.2) GERAÇÃO DO PDF COMPLETO (ReportLab)
        with col_exp2:
            output_pdf = BytesIO()
            doc = SimpleDocTemplate(output_pdf, pagesize=letter)
            styles = getSampleStyleSheet()
            elements = [
                Paragraph("<b>Relatório de Auditoria de Ativos - Metodologia DRC</b>", styles['Title']),
                Spacer(1, 15),
                Paragraph("Este documento apresenta a análise de conformidade de preços em relação aos bancos oficiais homologados.", styles['BodyText']),
                Spacer(1, 10)
            ]
            # Tabela simples de dados para anexar no PDF
            dados_pdf = [["Ativo", "Preço Forn.", "Preço Ref."]] + [[r['tipo_serviço_ativo'], f"R$ {r['preco_fornecedor']:.2f}", f"R$ {r['preco']:.2f}"] for _, r in df_total.iterrows()]
            t = Table(dados_pdf)
            t.setStyle(TableStyle([('BACKGROUND', (0,0), (-1,0), colors.grey), ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke), ('ALIGN', (0,0), (-1,-1), 'CENTER'), ('GRID', (0,0), (-1,-1), 1, colors.black)]))
            elements.append(t)
            doc.build(elements)
            
            st.download_button(
                label="📥 Baixar em PDF (.pdf)",
                data=output_pdf.getvalue(),
                file_name="Relatorio_Auditoria_Ativos.pdf",
                mime="application/pdf"
            )
            
        # 10.3) GERAÇÃO DO POWERPOINT (python-pptx)
        with col_exp3:
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5]) # Slide de título + Conteúdo
            slide.shapes.title.text = "Relatório Executivo de Engenharia e DRC"
            
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.text = "Destaques e Conclusões Técnicas da Análise:"
            for _, r in df_total.iterrows():
                p = tf.add_paragraph()
                p.text = f"• Ativo: {r['tipo_serviço_ativo']} | Desvio identificado em relação ao banco {r['banco']}."
                p.level = 1
                
            output_ppt = BytesIO()
            prs.save(output_ppt)
            st.download_button(
                label="📥 Baixar em PowerPoint (.pptx)",
                data=output_ppt.getvalue(),
                file_name="Apresentacao_Auditoria_Ativos.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
    else:
        st.warning("Insira dados de propostas e garanta que existam preços de referência cadastrados para liberar os Dashboards e as rotinas de exportação.")
